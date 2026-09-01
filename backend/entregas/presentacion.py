"""Genera la presentacion de PowerPoint de evidencia fotografica a partir
de la plantilla real (entregas/plantillas/Formato_rutas.pptx, copia de
docs/Formato_rutas.pptx). Ver el plan de esta funcion para el detalle de
por que se hace asi -- resumen: pptxgenjs (el generador anterior, en el
navegador) no puede abrir un .pptx existente, asi que esto se hace en el
backend con python-pptx, que si puede editar/clonar diapositivas reales
conservando el diseno (logos, fondo institucional, colores) tal cual.

Duplicar una diapositiva no es una funcion nativa de python-pptx -- se
hace copiando el arbol de formas de la diapositiva origen hacia una
diapositiva nueva (creada con el mismo layout, que ya trae el fondo/logos
solos por venir del layout) y volviendo a relacionar cualquier imagen que
se copie sin cambios: las relaciones (rId) son por diapositiva, copiar el
XML de una forma no copia sus relaciones.
"""

import copy
import io
import logging
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from . import storage
from .regiones import orden_regiones, region_de

logger = logging.getLogger(__name__)

RUTA_PLANTILLA = Path(__file__).resolve().parent / "plantillas" / "Formato_rutas.pptx"

_NS_REL = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# La plantilla original traia 8 rectangulos con foto de muestra (uno por
# slot) -- se quitaron porque solo eran de muestra y pesaban mucho. Ya no
# hay una forma que rellenar; se inserta una Picture nueva en la posicion
# exacta donde estaba cada rectangulo (medidas EMU tomadas de la plantilla
# original antes de que se quitaran, ver el plan). Cada posicion va
# emparejada con el id del cuadro de texto de su leyenda, que si sigue
# existiendo.
_ANCHO_FOTO = 2340000
_ALTO_FOTO = 2340000
_COLUMNAS_X = [874205, 3394298, 5914391, 8434484]
_FILAS_Y = [768534, 3832995]
_IDS_LEYENDA = [2, 25, 27, 29, 31, 33, 35, 37]
_POSICIONES_FOTO = [(x, y, _ANCHO_FOTO, _ALTO_FOTO) for y in _FILAS_Y for x in _COLUMNAS_X]
_SLOTS_FOTO = list(zip(_POSICIONES_FOTO, _IDS_LEYENDA))

_ID_TITULO_ENTIDAD = 9
_ID_DIA_CONTENIDO = 11
# Antes del ajuste que le hizo el usuario a la plantilla, la portada usaba
# un layout distinto y el texto de dia era el shape id=3. Ahora la
# portada comparte layout con la diapositiva de region y el id=3 quedo
# libre para el titulo -- el texto de dia se recorrio a id=4.
_ID_DIA_PORTADA = 4
_ID_TITULO_REGION = 2


def _texto_forma(forma, texto_nuevo):
    """Cambia el texto de una forma preservando el formato (fuente, tamano,
    color) del primer run -- reemplazar text_frame.text de un jalon lo
    resetea a la fuente por default."""
    parrafo = forma.text_frame.paragraphs[0]
    if not parrafo.runs:
        parrafo.add_run()
    parrafo.runs[0].text = texto_nuevo
    for run_extra in parrafo.runs[1:]:
        run_extra.text = ""


def _texto_leyenda(forma, nombre_unidad, clues):
    """La leyenda de cada foto son 2 parrafos separados en la plantilla
    (nombre de la unidad arriba, CLUES abajo) -- no una sola linea."""
    parrafos = forma.text_frame.paragraphs
    for parrafo, texto in zip(parrafos, [nombre_unidad, clues]):
        if not parrafo.runs:
            parrafo.add_run()
        parrafo.runs[0].text = texto
        for run_extra in parrafo.runs[1:]:
            run_extra.text = ""


def _forma_por_id(diapositiva, shape_id):
    for forma in diapositiva.shapes:
        if forma.shape_id == shape_id:
            return forma
    raise ValueError(f"No se encontro la forma id={shape_id} en la diapositiva -- "
                      "revisar si la plantilla Formato_rutas.pptx cambio de estructura.")


def _remapear_imagenes(elemento, parte_origen, parte_nueva):
    for nodo in elemento.iter():
        for atributo in ("embed", "link"):
            rid_viejo = nodo.get(f"{_NS_REL}{atributo}")
            if not rid_viejo:
                continue
            parte_imagen = parte_origen.related_part(rid_viejo)
            rid_nuevo = parte_nueva.relate_to(parte_imagen, RT.IMAGE)
            nodo.set(f"{_NS_REL}{atributo}", rid_nuevo)


def _eliminar_diapositiva(prs, diapositiva):
    """python-pptx no trae 'borrar diapositiva' -- hay que quitar su entrada
    de sldIdLst y soltar la relacion en el part de la presentacion, si no
    el archivo queda con una referencia rota."""
    for id_slide in list(prs.slides._sldIdLst):
        if prs.part.related_part(id_slide.rId) is diapositiva.part:
            prs.part.drop_rel(id_slide.rId)
            prs.slides._sldIdLst.remove(id_slide)
            return


def _duplicar_diapositiva(prs, diapositiva_origen):
    nueva = prs.slides.add_slide(diapositiva_origen.slide_layout)
    # add_slide() trae los placeholders vacios del layout -- no sirven,
    # se quitan antes de copiar encima las formas reales.
    for forma in list(nueva.shapes):
        forma._element.getparent().remove(forma._element)
    for forma_origen in diapositiva_origen.shapes:
        copia = copy.deepcopy(forma_origen._element)
        _remapear_imagenes(copia, diapositiva_origen.part, nueva.part)
        nueva.shapes._spTree.append(copia)
    return nueva


def _llenar_diapositiva_contenido(diapositiva, entidad_nombre, dia_etiqueta, fotos):
    _texto_forma(_forma_por_id(diapositiva, _ID_TITULO_ENTIDAD), entidad_nombre.title())
    _texto_forma(_forma_por_id(diapositiva, _ID_DIA_CONTENIDO), dia_etiqueta)

    for i, (posicion, id_texto) in enumerate(_SLOTS_FOTO):
        forma_texto = _forma_por_id(diapositiva, id_texto)
        agregada = False
        if i < len(fotos):
            visita = fotos[i]["visita"]
            evidencia = fotos[i]["evidencia"]
            try:
                bytes_imagen = storage.descargar_evidencia(evidencia.ruta_almacen)
                left, top, ancho, alto = posicion
                diapositiva.shapes.add_picture(io.BytesIO(bytes_imagen), left, top, ancho, alto)
                _texto_leyenda(forma_texto, visita.unidad_medica.nombre, visita.unidad_medica_id)
                agregada = True
            except Exception:
                # Una foto que no se puede insertar (formato que Pillow no
                # reconoce -- ej. HEIC subido antes de que se convirtiera
                # automaticamente a JPEG al subir, o un archivo corrupto) no
                # debe tumbar la presentacion completa -- se salta esa sola
                # y se deja constancia en los logs para poder darle
                # seguimiento (ver LOGGING en settings.py).
                logger.warning(
                    "No se pudo insertar la evidencia id=%s (unidad %s) en la presentacion",
                    evidencia.id, visita.unidad_medica_id, exc_info=True,
                )
        if not agregada:
            forma_texto._element.getparent().remove(forma_texto._element)


def construir_presentacion(dia_texto, fotos):
    """fotos: lista de {"visita": ProgramacionVisita, "evidencia": EvidenciaArchivo}
    (instancias de modelo reales -- ver views.py, ahi se valida y resuelve
    cada una desde la base antes de llegar aqui). Regresa un BytesIO listo
    para mandar como respuesta binaria."""
    prs = Presentation(str(RUTA_PLANTILLA))
    diapositiva_portada = prs.slides[0]
    diapositiva_region_base = prs.slides[1]
    diapositiva_contenido_base = prs.slides[2]

    dia_etiqueta = f"Día 1: {dia_texto}"
    _texto_forma(_forma_por_id(diapositiva_portada, _ID_DIA_PORTADA), dia_etiqueta)

    # Region (orden fijo) -> entidad (orden alfabetico) -> fotos de esa entidad.
    por_region = {}
    for foto in fotos:
        entidad_nombre = foto["visita"].unidad_medica.entidad.nombre
        region = region_de(entidad_nombre)
        por_region.setdefault(region, {}).setdefault(entidad_nombre, []).append(foto)

    # diapositiva_region_base y diapositiva_contenido_base son PLANTILLAS --
    # nunca se editan/mutan directamente (llenar una diapositiva de
    # contenido con menos de 8 fotos borra los slots sobrantes; si eso le
    # pasara a la base, la siguiente copia saldria incompleta). Cada
    # region/entidad, incluida la primera, usa una copia nueva; al final se
    # borran las 2 plantillas, que nunca deben quedar en el resultado.
    for region in orden_regiones():
        entidades = por_region.get(region)
        if not entidades:
            continue

        diapositiva_region = _duplicar_diapositiva(prs, diapositiva_region_base)
        _texto_forma(_forma_por_id(diapositiva_region, _ID_TITULO_REGION), region)

        for entidad_nombre in sorted(entidades):
            fotos_entidad = entidades[entidad_nombre]
            for inicio in range(0, len(fotos_entidad), len(_SLOTS_FOTO)):
                grupo = fotos_entidad[inicio: inicio + len(_SLOTS_FOTO)]
                diapositiva_contenido = _duplicar_diapositiva(prs, diapositiva_contenido_base)
                _llenar_diapositiva_contenido(diapositiva_contenido, entidad_nombre, dia_etiqueta, grupo)

    _eliminar_diapositiva(prs, diapositiva_region_base)
    _eliminar_diapositiva(prs, diapositiva_contenido_base)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
